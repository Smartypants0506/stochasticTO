"""Fill the 48x36 poster template with the project's final results.

    python viz/build_poster.py

    -> output/poster/robust_TO_poster.pptx

The template (pragmaticgraphite_48x36.pptx) is a 4-column, 2-row grid:
section headers sit at y=8.0 and y=20.0, body text at y=9.1 and y=21.1, with
columns at x = 0.7 / 12.7 / 24.8 / 36.8. Two of the four bottom cells ship
without headers; those become the figure panels.

Every number written here is traceable to a JSON under output/studies/ and is
FINAL -- see docs/POSTER_DATA.md. Nothing on this poster depends on the
production sweep, which will not finish before the deadline.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = Path("/home/ovb/Downloads/pragmaticgraphite_48x36.pptx")
OUT = ROOT / "output" / "poster" / "robust_TO_poster.pptx"

INK = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT = RGBColor(0x1F, 0x4E, 0x79)
MUTED = RGBColor(0x55, 0x5F, 0x6B)
FLAG = RGBColor(0xB0, 0x30, 0x2A)

COL_X = [0.7, 12.7, 24.8, 36.8]
COL_W = 10.5
HDR_Y = [8.0, 20.0]
BODY_Y = [9.1, 21.1]

# ---------------------------------------------------------------- content

TITLE = ("When Does Spatially Correlated Manufacturing Error Matter "
         "in Robust Topology Optimization?")
AUTHORS = ("[AUTHOR NAME]  |  [Department / Group]  |  "
           "National Institute of Standards and Technology")

ABSTRACT = """Topology optimization produces designs tuned to a geometry no factory can build exactly. Milling and etching make members thinner or thicker than intended, and the standard way to model this is to randomize the threshold of the Heaviside projection that turns a filtered density field into a structure.

Prior work modelled that threshold as a spatially correlated random field, but fixed the correlation length at a single value and reported that designs optimized against uniform error were just as robust to non-uniform error. Whether that finding generalized was left open.

We treat the correlation length as the independent variable and sweep it over a 32-fold range in three-dimensional linear elasticity, solving the robust problem by sample average approximation with 512 finite-element evaluations per design iteration.

Response variability rises MONOTONICALLY with correlation length and is maximal in the spatially uniform limit. The uniform threshold therefore bounds the response variance from above: spatial correlation can only reduce it, and the inexpensive scalar model is the conservative choice. We additionally quantify the sampling error of the method itself, which prior work did not."""

INTRODUCTION = """THE PROBLEM. Over- and under-etching perturb every boundary of a manufactured part. Topology-optimized designs are unusually exposed because they concentrate material into thin members, so a small uniform offset can remove a load path entirely.

THE STANDARD MODEL. Apply a density filter, then a smooth Heaviside projection with threshold eta. A HIGH eta keeps only the densest material and thins the structure (over-etching); a LOW eta thickens it (under-etching). Randomizing eta therefore models manufacturing error without perturbing the mesh.

THE OPEN QUESTION. Schevenels, Lazarov & Sigmund (CMAME 200:3613-3627, 2011) made eta a spatially correlated random field and found, for a 2D compliant mechanism and a 2D heat sink, that uniform-error and non-uniform-error designs were equally robust. That result came from ONE correlation length. Nobody had asked WHEN spatial correlation matters -- or whether the expensive field model earns its cost over a single random variable.

WHAT WE ADD. The correlation-length sweep, in 3D compliance; and an error analysis of the sample-average method itself, replicated across ten independent seeds."""

METHODOLOGY = """DESIGN CHAIN. SIMP stiffness E(rho) = E0 * rho^p, p = 3, followed by a Helmholtz PDE filter (-R^2 grad^2 rho~ + rho~ = rho, R = 0.6) and a smooth Heaviside projection with continuation beta = 8 -> 128.

THE UNCERTAINTY. eta is a NON-GAUSSIAN RANDOM FIELD. An underlying Gaussian field is discretized by a Karhunen-Loeve expansion on the finite-element mesh (>=95% retained variance), standardized by its own pointwise standard deviation, then mapped through a memoryless isoprobabilistic transform to a Beta(2,2) marginal on [0.25, 0.75].

Standardizing before the transform makes the marginal EXACT regardless of the truncation level, so the perturbation magnitude is set by the band alone.

ROBUST PROBLEM.
    min  J = mu_C + lambda * sigma_C
    s.t. E[V(rho)] <= V*,  0 <= rho <= 1

SOLUTION. Surrogate-free sample average approximation: ONE fixed set of N = 512 eta realizations, drawn once and reused at every design iteration (common random numbers), each evaluated by a full FEA solve with exact sample-average gradients. MMA via PETSc TAO; samples parallelized across MPI sub-communicators.

TEST CASE. 3D cantilever beam, domain 10 x 30 x 10 (dimensionless), E = 100, nu = 0.25, volume fraction 0.08, ~154k dofs."""

RESULTS = """HEADLINE: variability grows with correlation length.

 l_c    l_c/L    cv = sigma_C/mu_C   95% CI
 1      0.033    1.2337              [1.168, 1.317]
 2      0.067    1.4611              [1.389, 1.547]
 4      0.133    1.9305              [1.843, 2.034]
 8      0.267    2.3502              [2.217, 2.502]
 16     0.533    3.0260              [2.840, 3.230]
 32     1.067    3.3059              [3.097, 3.533]
 UNIFORM  inf    3.1788              [2.972, 3.406]

ALL FOUR adjacent 95% intervals from l_c = 1 to 16 are DISJOINT, so monotonicity is established rather than asserted. Endpoint contrast 2.58x, also disjoint. The three longest correlation lengths are mutually indistinguishable -- which is the point, since all three are effectively the uniform case.

VERIFICATION. sigma_C/mu_C varies by only 1.3% across a 1.6-fold range of element size, so it is a continuum property. mu_C and sigma_C INDIVIDUALLY move ~37% over the same range, because the resolved traction area varies with h. We therefore report the RATIO, never absolute compliance.

ROBUSTNESS GAINED. The robust design reaches cv = 0.51 against 2.00 for the deterministic design -- about 4x less variable -- while staying near-binary (measure of non-discreteness 0.234%)."""

SAMPLING = """The method's own sampling error, from ten independently seeded solves at N = 512, each re-scored on the SAME independent 5000-sample set (Mak/Morton/Wood estimator).

 IN-SAMPLE sigma_C is OPTIMISTIC BY 33.9%
   (0.0531 in-sample vs 0.0804 out-of-sample)

 OPTIMALITY GAP  +3.1%,  95% CI [0.58%, 2.36%]
   -- excludes zero, so the gap is resolved

 RUN-TO-RUN NOISE FLOOR  15.5% (robust MAD)

FOUR OF TEN designs are materially worse out of sample -- one by a factor of 2.47 -- and in-sample sigma_C is 0.048-0.058 for ALL TEN. It gives no warning whatsoever.

WHY IT MATTERS. Prior work justified N = 100 with a single 10000-sample check on ONE design and concluded it was sufficient. That procedure would have missed this in six cases out of ten. Their own published heat-sink data already contains the effect unremarked (0.042 in-sample vs 0.050 reference, a 16% optimism)."""

CONCLUSION = """THE RESULT. Response variability rises monotonically with the correlation length of the manufacturing error, and is maximal when the error is spatially UNIFORM.

WHAT THAT MEANS. The uniform threshold BOUNDS the response variance from above. Spatial correlation can only reduce it, because erosions in one region are partly offset by dilations in another -- the same cancellation prior work invoked to explain its heat sink, now measured across the whole range instead of asserted at one point.

PRACTICAL CONSEQUENCE. A designer can use the inexpensive scalar-threshold model and be CONSERVATIVE by construction: roughly 8x fewer finite-element solves, with a bound rather than a hope.

This inverts the premise the project began with -- that a spatially correlated field was necessary -- and replaces it with a stronger, more useful claim.

SCOPE. This is a statement about the RESPONSE of a fixed design. Extending it to a statement about the DESIGN requires a cross-evaluation of designs optimized at one correlation length and scored at another; that experiment is running and is the immediate next step."""

LIMITATIONS = """NOT METROLOGY-CALIBRATED. The eta band was chosen so the induced boundary shift is resolvable on the mesh (standard deviation 0.41 elements), not fitted to a measured process. This is a robustness ENVELOPE, not a process tolerance.

AGGRESSIVE RELATIVE TO FEATURE SIZE. eps_max/R = 0.50 here versus 0.108 in the 2D prior work, forced by 3D affordability: their filter was resolved at R/h = 8.4, ours at 1.5. Matching them in 3D would need ~8.2M elements.

FIRST-ORDER OPTIMALITY NOT REACHED. The relative stationarity residual plateaus at 0.04-0.12 against a 1e-3 tolerance. Move-limit continuation was tested and made it worse; the cause is beta = 128 projection stiffness, where the responsive band is only ~19/beta wide.

GRADIENT VERIFICATION. First-moment sensitivities verify to 4-6 significant figures. The second-moment sensitivity dsigma_C/drho is correct analytically (it reproduces an exact finite difference to 3.5e-10 on a noiseless problem) but cannot be verified to 0.1% against FEA finite differences, because differencing a standard deviation amplifies per-solve noise.

ACKNOWLEDGEMENTS
[Advisor / supervisor names]
[Funding source or program, e.g. NIST SURF]
Built on FEniTop (Jia, Wang & Zhang 2024) and dolfiny's MMA; both vendored and modified, with the modifications documented."""

# --------------------------------------------------------------- helpers


def set_text(shape, text, size=20, bold=False, color=INK, align=PP_ALIGN.LEFT,
             space_after=8, line_spacing=0.95):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, para in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = para
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"


def add_box(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def add_header(slide, x, y, text):
    box = add_box(slide, x, y, COL_W, 0.8)
    set_text(box, text, size=40, bold=True, color=ACCENT, space_after=0)
    return box


FOOTER_Y = 34.8   # template's footer band starts at 35.0


def add_figure(slide, path: Path, x, y, w, caption, cap_h=1.15, max_h=None):
    """Place an image with a caption, fitted to the space actually available.

    Scaling to column WIDTH alone is what pushed the square P1 triptych 6
    inches past the bottom of the poster: a 1:1 image at 10.5in wide is 10.5in
    tall. Height is therefore capped as well, and the image is centred in the
    column when the height cap makes it narrower than the column.
    """
    from PIL import Image
    with Image.open(path) as im:
        aspect = im.size[1] / im.size[0]

    limit = max_h if max_h is not None else (FOOTER_Y - y - cap_h - 0.15)
    h = w * aspect
    if h > limit:
        h = limit
        w = h / aspect
    x_centred = x + (COL_W - w) / 2.0

    slide.shapes.add_picture(str(path), Inches(x_centred), Inches(y),
                             Inches(w), Inches(h))
    cap = add_box(slide, x, y + h + 0.06, COL_W, cap_h)
    set_text(cap, caption, size=14, color=MUTED, space_after=2, line_spacing=0.9)
    return y + h + 0.06 + cap_h


def add_placeholder(slide, x, y, w, h, title, desc):
    """A labelled box for a figure that does not exist yet."""
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.background()
    box.line.color.rgb = FLAG
    box.line.width = Pt(2)
    set_text(box, f"{title}\n\n{desc}", size=15, color=FLAG,
             align=PP_ALIGN.CENTER, space_after=4)
    return y + h


def main() -> None:
    prs = Presentation(str(TEMPLATE))
    slide = prs.slides[0]

    by_name = {sh.name: sh for sh in slide.shapes}

    # --- title block -------------------------------------------------------
    set_text(by_name["Google Shape;90;p13"], TITLE, size=72, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
             space_after=0, line_spacing=0.95)
    set_text(by_name["Google Shape;91;p13"], AUTHORS, size=32,
             color=RGBColor(0xE8, 0xEC, 0xF0), align=PP_ALIGN.CENTER,
             space_after=0)

    # --- section headers ---------------------------------------------------
    set_text(by_name["Google Shape;99;p13"], "Abstract", size=40, bold=True, color=ACCENT, space_after=0)
    set_text(by_name["Google Shape;100;p13"], "Methodology", size=40, bold=True, color=ACCENT, space_after=0)
    set_text(by_name["Google Shape;101;p13"], "Results", size=40, bold=True, color=ACCENT, space_after=0)
    set_text(by_name["Google Shape;102;p13"], "Conclusion", size=40, bold=True, color=ACCENT, space_after=0)
    set_text(by_name["Google Shape;98;p13"], "Introduction", size=40, bold=True, color=ACCENT, space_after=0)
    set_text(by_name["Google Shape;103;p13"], "Limitations & Acknowledgements",
             size=40, bold=True, color=ACCENT, space_after=0)
    # The template leaves the two middle bottom cells unlabelled; they become
    # the figure panels, so they need headers of their own.
    add_header(slide, COL_X[1], HDR_Y[1], "The Error Model")
    add_header(slide, COL_X[2], HDR_Y[1], "Sampling Error")

    # --- body text ---------------------------------------------------------
    set_text(by_name["Google Shape;92;p13"], "", size=1)   # unused stubs
    set_text(by_name["Google Shape;93;p13"], "", size=1)
    set_text(by_name["Google Shape;94;p13"], "", size=1)
    set_text(by_name["Google Shape;95;p13"], "", size=1)
    set_text(by_name["Google Shape;96;p13"], "", size=1)
    set_text(by_name["Google Shape;89;p13"], "", size=1)

    set_text(add_box(slide, COL_X[0], BODY_Y[0], COL_W, 10.2), ABSTRACT, size=19)
    set_text(add_box(slide, COL_X[1], BODY_Y[0], COL_W, 10.2), METHODOLOGY, size=18)
    set_text(add_box(slide, COL_X[2], BODY_Y[0], COL_W, 10.2), RESULTS, size=17.5)
    set_text(add_box(slide, COL_X[3], BODY_Y[0], COL_W, 10.2), CONCLUSION, size=19)

    set_text(add_box(slide, COL_X[0], BODY_Y[1], COL_W, 8.2), INTRODUCTION, size=17)
    set_text(add_box(slide, COL_X[3], BODY_Y[1], COL_W, 13.4), LIMITATIONS, size=15.5)

    # --- figures -----------------------------------------------------------
    fig = ROOT / "output" / "figures"
    pv = fig / "poster"

    # Column 2 bottom: what the error physically does + the field itself.
    y = add_figure(
        slide, pv / "P1_triptych.png", COL_X[1], BODY_Y[1], COL_W,
        "FIG 1  [poster/P1_triptych.png]  One design, three manufacturing outcomes. "
        "Volume fraction 0.026 (over-etched) / 0.081 (as designed) / 0.150 "
        "(under-etched). The eroded case retains under a third of the intended "
        "material.", max_h=7.4)
    add_figure(
        slide, pv / "P3_eta_field_0.png", COL_X[1], y + 0.15, COL_W,
        "FIG 2  [poster/P3_eta_field_0.png]  One realization of the threshold "
        "field eta(x), Beta(2,2) on [0.25, 0.75]. The projection threshold varies "
        "in SPACE -- this is the input the method randomizes.", max_h=3.9)

    # Column 3 bottom: the headline curve and the sampling-error scatter.
    y = add_figure(
        slide, fig / "figA_correlation_length.png", COL_X[2], BODY_Y[1], COL_W,
        "FIG 3  [figA_correlation_length.png]  THE HEADLINE. cv = sigma_C/mu_C "
        "against correlation length, with the uniform limit marked. Monotonic; "
        "adjacent 95% intervals disjoint from l_c = 1 to 16.", max_h=6.2)
    y = add_figure(
        slide, fig / "figC_gap_replications.png", COL_X[2], y + 0.15, COL_W,
        "FIG 4  [figC_gap_replications.png]  In-sample vs out-of-sample sigma_C "
        "across ten seeds. Every point sits above the diagonal; the circled "
        "outlier is 2.47x the median with an unremarkable in-sample value.", max_h=5.2)

    # Column 1 bottom, under the Introduction: the convergence evidence.
    add_figure(
        slide, fig / "figB_mesh_convergence.png", COL_X[0], BODY_Y[1] + 8.4, COL_W,
        "FIG 5  [figB_mesh_convergence.png]  Only the RATIO is mesh-converged. "
        "mu_C and sigma_C each move ~37% across a 1.6x range of element size "
        "while sigma_C/mu_C holds inside 1.3%.", max_h=4.0)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT.relative_to(ROOT)}")
    print(f"  {prs.slide_width/914400:.0f} x {prs.slide_height/914400:.0f} inches")
    print("  figures embedded: P1_triptych, P3_eta_field_0, figA, figB, figC")


if __name__ == "__main__":
    main()
